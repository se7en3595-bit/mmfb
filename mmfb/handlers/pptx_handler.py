"""Powerpoint .pptx 格式处理器

职责：
1. 解析 .pptx 文件（幻灯片结构/文本/图片/形状），输出预览数据
2. 预览模式：前端渲染缩略图列表 + 当前幻灯片预览
3. 支持幻灯片切换、全屏演示模式、备注显示

技术方案：
- 用 python-pptx 解析 pptx 内部结构
- 将幻灯片（形状/文本/图片/位置）转换为结构化数据
- 图片以 base64 内嵌
- 前端根据EMU坐标按比例渲染幻灯片内容
"""
import base64
import os
from typing import Any, Dict, List, Optional

from mmfb.core.handler_base import BaseHandler

# pptx 内部坐标单位：EMU (English Metric Unit, 1 inch = 914400 EMU)
# 常见幻灯片尺寸：宽屏 16:9 = 9144000 x 5143500 EMU

PPTX_EXTENSIONS: List[str] = [
    ".pptx",
    ".pptm",
    ".potx",
    ".potm",
    ".ppsx",
    ".ppsm",
]


class PptxHandler(BaseHandler):
    """PowerPoint .pptx 文件处理器

    支持的扩展名：.pptx / .pptm / .potx / .potm / .ppsx / .ppsm

    预览输出结构：
    - mime:
        'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    - template: 'pptx'
    - data.slides[i].number: 幻灯片编号 (1-indexed)
    - data.slides[i].layoutName: 布局名称
    - data.slides[i].width / height: 幻灯片 EMU 宽高
    - data.slides[i].shapes[i].id / text / runs / left / top / width / height
    - data.slides[i].images[i].id / mime / base64 / left / top / width / height
    - data.slides[i].notes: 备注文本
    - data.slideCount: 总数
    - data.slideWidth / slideHeight: 统一EMU宽高
    - data.filePath / data.fileSize
    - editable: False
    """

    extensions = PPTX_EXTENSIONS

    def get_preview(self) -> Optional[Dict[str, Any]]:
        """获取 pptx 预览数据"""
        try:
            if not os.path.isfile(self.path):
                return self._make_error_result("file not found")

            file_size = os.path.getsize(self.path)

            try:
                from pptx import Presentation
                prs = Presentation(self.path)
            except Exception as e:
                return self._make_error_result(f"failed to open pptx: {e}", file_size=file_size)

            slide_width = prs.slide_width
            slide_height = prs.slide_height

            slides: List[Dict[str, Any]] = []

            for slide_idx, slide in enumerate(prs.slides):
                slide_data = self._extract_slide(slide, slide_idx, slide_width, slide_height)
                slides.append(slide_data)

            file_name = os.path.basename(self.path)

            return {
                "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "template": "pptx",
                "data": {
                    "slides": slides,
                    "slideCount": len(slides),
                    "slideWidth": slide_width,
                    "slideHeight": slide_height,
                    "file_path": os.path.abspath(self.path),
                    "file_name": file_name,
                    "file_size": file_size,
                },
                "editable": False,
            }
        except Exception as e:
            return self._make_error_result(str(e))

    def get_edit(self) -> None:
        """pptx v1 不支持就地编辑"""
        return None

    def _extract_slide(self, slide, slide_idx: int, prs_width: int, prs_height: int) -> Dict[str, Any]:
        """提取单页幻灯片数据"""
        # 布局名称
        layout_name = ""
        try:
            if slide.slide_layout and slide.slide_layout.name:
                layout_name = slide.slide_layout.name
        except Exception:
            layout_name = "Unknown"

        shapes_data: List[Dict[str, Any]] = []
        images_data: List[Dict[str, Any]] = []

        for shape in slide.shapes:
            shape_data = self._extract_shape(shape)
            if shape_data is None:
                continue
            if shape_data.get("type") == "image":
                images_data.append(shape_data)
            else:
                shapes_data.append(shape_data)

        # 备注
        notes_text = self._extract_notes(slide)

        return {
            "number": slide_idx + 1,
            "layoutName": layout_name,
            "width": prs_width,
            "height": prs_height,
            "shapes": shapes_data,
            "images": images_data,
            "notes": notes_text,
        }

    def _extract_shape(self, shape) -> Optional[Dict[str, Any]]:
        """提取单个形状数据"""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            shape_info: Dict[str, Any] = {
                "id": f"shape_{id(shape)}",
                "name": shape.name,
                "shapeType": str(shape.shape_type),
            }

            # 位置与尺寸（EMU）
            try:
                shape_info["left"] = int(shape.left or 0)
                shape_info["top"] = int(shape.top or 0)
                shape_info["width"] = int(shape.width or 0)
                shape_info["height"] = int(shape.height or 0)
            except Exception:
                shape_info["left"] = 0
                shape_info["top"] = 0
                shape_info["width"] = 0
                shape_info["height"] = 0

            # 图片类型
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return self._extract_picture(shape, shape_info)

            # 表格类型
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return self._extract_table(shape, shape_info)

            # 有文本框
            if shape.has_text_frame:
                paragraphs: List[Dict[str, Any]] = []
                full_text_parts: List[str] = []

                for para in shape.text_frame.paragraphs:
                    runs_data, runs_text = self._extract_paragraph_runs(para)

                    if runs_text:
                        full_text_parts.append(runs_text)

                    paragraphs.append({
                        "text": runs_text,
                        "level": getattr(para, 'level', 0) or 0,
                        "runs": runs_data,
                    })

                shape_info["type"] = "text"
                shape_info["text"] = "\n".join(full_text_parts)
                shape_info["paragraphs"] = paragraphs
                return shape_info

            # 无文本无图片的纯形状
            shape_info["type"] = "shape"
            shape_info["text"] = ""
            shape_info["paragraphs"] = []
            return shape_info

        except Exception as e:
            return {
                "id": f"shape_{id(shape)}",
                "name": getattr(shape, 'name', 'unknown'),
                "type": "error",
                "text": f"[parse error: {e}]",
                "paragraphs": [],
                "left": 0, "top": 0, "width": 0, "height": 0,
                "shapeType": "unknown",
            }

    def _extract_paragraph_runs(self, paragraph):
        """提取段落中所有 runs 文本与样式，返回 (runs_list, full_text)"""
        runs: List[Dict[str, Any]] = []
        full_text = ""

        try:
            para_runs = paragraph.runs
        except Exception:
            para_runs = []

        for run in para_runs:
            try:
                text = run.text or ""
            except Exception:
                text = ""

            if not text:
                continue

            run_info: Dict[str, Any] = {"text": text}

            # 字体样式
            try:
                font = run.font
                if font.bold:
                    run_info["bold"] = True
                if font.italic:
                    run_info["italic"] = True
                if font.underline:
                    run_info["underline"] = True
                if font.size:
                    # font.size 以 EMU 为单位存储，但值就是磅数对应的 EMU
                    # 例如 12pt = 12 * 12700 = 152400 EMU
                    # 转回 pt
                    try:
                        size_emu = int(font.size)
                        if size_emu > 0:
                            run_info["sizePt"] = round(size_emu / 12700, 1)
                    except Exception:
                        pass
                if font.color and font.color.rgb is not None:
                    run_info["color"] = str(font.color.rgb)
                if font.name:
                    run_info["fontName"] = font.name
            except Exception:
                pass

            runs.append(run_info)
            full_text += text

        return runs, full_text

    def _extract_picture(self, shape, shape_info: Dict[str, Any]) -> Dict[str, Any]:
        """提取图片形状"""
        shape_info["type"] = "image"
        shape_info["text"] = ""
        shape_info["paragraphs"] = []

        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                blob = image.blob
                content_type = image.content_type or "image/png"
                b64 = base64.b64encode(blob).decode("ascii")
                shape_info["imageData"] = b64
                shape_info["imageMime"] = content_type
            else:
                shape_info["imageData"] = ""
                shape_info["imageMime"] = ""
        except Exception as e:
            shape_info["imageData"] = ""
            shape_info["imageMime"] = ""
            shape_info["_imageError"] = str(e)

        return shape_info

    def _extract_table(self, shape, shape_info: Dict[str, Any]) -> Dict[str, Any]:
        """提取表格形状为结构化数据"""
        shape_info["type"] = "table"
        shape_info["text"] = ""
        shape_info["paragraphs"] = []

        rows: List[List[str]] = []
        try:
            table = shape.table
            for row in table.rows:
                cells: List[str] = []
                for cell in row.cells:
                    cell_text = (cell.text or "").strip()
                    cells.append(cell_text)
                rows.append(cells)
        except Exception:
            pass

        shape_info["rows"] = rows
        return shape_info

    def _extract_notes(self, slide) -> str:
        """提取幻灯片的备注文本"""
        try:
            if not slide.has_notes_slide:
                return ""
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            parts: List[str] = []
            for para in text_frame.paragraphs:
                para_text = ""
                for run in para.runs:
                    para_text += run.text or ""
                if para_text:
                    parts.append(para_text)
            return "\n".join(parts)
        except Exception:
            return ""

    def _make_error_result(self, error: str, file_size: int = 0) -> Dict[str, Any]:
        """构造错误结果"""
        return {
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "template": "pptx",
            "data": {
                "slides": [],
                "slideCount": 0,
                "slideWidth": 9144000,
                "slideHeight": 5143500,
                "file_path": os.path.abspath(self.path),
                "file_name": os.path.basename(self.path),
                "file_size": file_size,
            },
            "editable": False,
            "error": error,
        }
